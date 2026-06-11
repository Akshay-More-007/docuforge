"""
ppt_builder.py — Build a .pptx from structured content.
Accepts a list of slide dicts from the LLM and creates a presentation.

Two modes:
  1. template_path=None → blank Presentation (original approach)
  2. template_path=source.pptx → open source, modify slides IN PLACE
     This preserves ALL design elements (group shapes, gradient rectangles, logos,
     CGI theme, Aptos fonts) while updating only the text/table content.

Each slide dict the LLM returns:
{
    "layout": "title" | "title_content" | "two_column" | "blank",
    "title": "...",
    "content": "string or list of bullet strings",
    "table": {                   # use instead of content for data slides
        "headers": ["Col1", "Col2"],
        "rows": [["R1C1", "R1C2"], ["R2C1", "R2C2"]]
    },
    "left_content": "...",       # for two_column
    "right_content": "...",      # for two_column
    "notes": "...",              # optional speaker notes
    "theme": {                   # used for fallback/blank mode
        "fonts": ["Primary Font", "Body Font"],
        "font_sizes": [title_pt, body_pt],
        "text_colors": ["RRGGBB_title", "RRGGBB_body"],
        "background_color": "RRGGBB or null"
    }
}

CGI table style applied when template is provided:
    {5C22544A-7EE6-4342-B048-85BDC9FD1C3A}  firstRow=1, bandRow=1
    No explicit cell fills — theme handles header color + alternating rows.
"""

import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# CGI table style GUID — produces dark header + banded rows from theme colors
CGI_TABLE_STYLE_ID = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"

# "No Style, No Grid" — transparent table used with a gradient header bar behind
# the first row (the CGI clean-deck design).
TRANSPARENT_TABLE_STYLE_ID = "{5940675A-B579-460E-94D1-54222C63F5DA}"

# CGI brand gradient (red → magenta → purple) for the header bar fallback
CGI_GRADIENT_STOPS = [(0, "E31937"), (60000, "A82465"), (100000, "5236AB")]

# Uniform table-slide geometry (matches the CGI reference design)
REF_GEOM = {
    "left":   Inches(0.53),
    "top":    Inches(0.79),
    "width":  Inches(12.02),
    "bar_h":  Inches(0.54),
    "title_top": Inches(0.25),
    "title_h":   Inches(0.54),
    "bottom_margin": Inches(0.15),
}


# ── Public entry point ────────────────────────────────────────────────────────

def build_pptx(slides: list[dict], output_path: str, template_path: str | None = None) -> str:
    """
    Build a .pptx from slide data.

    If template_path is provided (the source PPTX), slides are modified IN PLACE
    so all design elements (GROUP shapes, gradient rects, logos, CGI theme) are
    preserved. Content (title text, tables, agenda items) is updated.

    If template_path is None, a fresh blank presentation is built (fallback).

    Returns: absolute path to saved .pptx
    """
    if template_path:
        out = _build_pptx_from_template(slides, output_path, template_path)
    else:
        out = _build_pptx_blank(slides, output_path)
    return out


def _build_pptx_from_template(slides: list[dict], output_path: str, template_path: str) -> str:
    """
    Open source PPTX as template and modify each slide IN PLACE.
    Preserves design elements; updates text/table content only.
    """
    prs = Presentation(template_path)
    source_slides = list(prs.slides)
    n_source = len(source_slides)
    n_target = len(slides)

    logger.info(
        f"[PPTBuilder] Template mode: {n_source} source slides → {n_target} output slides"
    )

    # Reference design: a source slide that already pairs a table with a gradient
    # header bar defines the uniform look every other table slide is rebuilt to.
    ref = _find_reference_design(prs)

    # ── Modify existing slides ────────────────────────────────────────────────
    for i, slide_data in enumerate(slides):
        if i < n_source:
            _modify_slide_in_place(source_slides[i], slide_data, ref)
        else:
            # More output slides than source → add fresh slides
            layout_map = _build_layout_map(prs)
            theme = slide_data.get("theme") or {}
            _add_new_slide_to_prs(prs, layout_map, slide_data, theme)

    # ── Remove surplus source slides (from the end) ───────────────────────────
    while len(prs.slides) > n_target:
        _remove_last_slide(prs)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    logger.info(f"[PPTBuilder] Saved pptx (template mode) → {out}")
    return str(out.resolve())


def _build_pptx_blank(slides: list[dict], output_path: str) -> str:
    """Build a fresh blank presentation (no template). Original approach."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    layout_map = _build_layout_map(prs)

    for slide_data in slides:
        theme = slide_data.get("theme") or {}
        _add_new_slide_to_prs(prs, layout_map, slide_data, theme)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    logger.info(f"[PPTBuilder] Saved pptx (blank mode) → {out}")
    return str(out.resolve())


# ── In-place slide modification (template mode) ───────────────────────────────

def _modify_slide_in_place(slide, data: dict, ref: dict | None = None) -> None:
    """
    Update a slide's content while preserving its design elements.

    Modifies:
    - Title placeholder text (run formatting preserved)
    - For table slides: normalizes the slide to the deck's reference design —
      uniform title strip, gradient header bar, transparent table style
    - For agenda slides: fills numbered placeholder text
    - For cover slide: updates subtitle TextBox if present

    Preserves:
    - GROUP shapes (design graphics)
    - PICTURE shapes (logos)
    - Existing gradient header bars (reused, never duplicated)
    - Background and layout
    """
    title_text = data.get("title", "")
    table_data = data.get("table")
    content    = data.get("content", "")
    slide_num  = data.get("slide_number", 0)

    is_agenda = _is_agenda_slide(data)

    # ── 1. Update title placeholder (ph[0]) ───────────────────────────────────
    if title_text:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                _replace_placeholder_text(ph, title_text)
                # Table slides: normalize the title strip geometry (keep the
                # slide's own left/width so the design stays untouched)
                if table_data:
                    ph.top    = REF_GEOM["title_top"]
                    ph.height = REF_GEOM["title_h"]
                break

    # ── 2. Table slides: normalize to the reference design ────────────────────
    if table_data:
        existing_bar = _slide_gradient_bar(slide)
        existing_tbl = next(
            (s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE), None
        )

        # Slide already in reference shape (e.g. the reference slide itself) and
        # the table dimensions are unchanged → just update cell texts in place.
        updated_in_place = (
            existing_bar is not None
            and existing_tbl is not None
            and _update_table_texts_in_place(existing_tbl, table_data)
        )
        if not updated_in_place:
            _remove_table_shapes(slide)
            if existing_bar is None:
                _add_header_bar(slide, ref)
            else:
                existing_bar.left   = REF_GEOM["left"]
                existing_bar.top    = REF_GEOM["top"]
                existing_bar.width  = REF_GEOM["width"]
                existing_bar.height = REF_GEOM["bar_h"]
            _add_reference_table(slide, table_data, ref)

    # ── 3. Agenda slide: fill numbered section placeholders ───────────────────
    elif is_agenda:
        _fill_agenda_placeholders(slide, content)

    # ── 4. Cover slide (slide 1): update subtitle TextBox ─────────────────────
    elif slide_num == 1:
        _update_cover_subtitle(slide, content)

    # ── 5. Speaker notes ──────────────────────────────────────────────────────
    notes_text = data.get("notes", "")
    if notes_text:
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
        except Exception as e:
            logger.warning(f"[PPTBuilder] Notes error: {e}")


def _replace_placeholder_text(ph, text: str) -> None:
    """Replace all text in a placeholder, preserving the first run's formatting."""
    tf = ph.text_frame
    # Clear all paragraphs except the first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    para = tf.paragraphs[0]
    runs = para.runs
    if runs:
        # Reuse the first run (keeps its rPr: size, color, font) — drop the rest
        runs[0].text = text
        for run in runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        run = para.add_run()
        run.text = text


# ── Reference design (gradient header bar + transparent table) ────────────────

def _slide_gradient_bar(slide):
    """Find a slim, wide AUTO_SHAPE with a gradient fill (a header bar)."""
    for shape in slide.shapes:
        try:
            if shape.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE,):
                continue
            if shape.height < Inches(1.0) and shape.width > Inches(6.0) \
                    and "gradFill" in shape._element.xml:
                return shape
        except Exception:
            continue
    return None


def _find_reference_design(prs) -> dict:
    """
    Scan the deck for a slide that pairs a TABLE with a gradient header bar —
    that slide defines the clean design every table slide is normalized to.
    Falls back to the CGI defaults when no such slide exists.
    """
    ref = {
        "bar_el": None,
        "style_id": TRANSPARENT_TABLE_STYLE_ID,
        "font": "Arial",
        "hdr_sz": 18,
        "body_sz": 11,
    }
    import re as _re
    for slide in prs.slides:
        bar = _slide_gradient_bar(slide)
        tbl_shape = next(
            (s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE), None
        )
        if bar is None or tbl_shape is None:
            continue
        ref["bar_el"] = bar._element
        xml = tbl_shape._element.xml
        m = _re.search(r"<a:tableStyleId>([^<]+)</a:tableStyleId>", xml)
        if m:
            ref["style_id"] = m.group(1)
        m = _re.search(r'typeface="([^"]+)"', xml)
        if m:
            ref["font"] = m.group(1)
        trs = _re.findall(r"<a:tr .*?</a:tr>", xml, _re.DOTALL)
        if trs:
            szs = _re.findall(r'sz="(\d+)"', trs[0])
            if szs:
                ref["hdr_sz"] = int(szs[0]) / 100
        if len(trs) > 1:
            szs = _re.findall(r'sz="(\d+)"', trs[1])
            if szs:
                ref["body_sz"] = int(szs[0]) / 100
        logger.info(
            f"[PPTBuilder] Reference design found: style={ref['style_id']} "
            f"font={ref['font']} hdr={ref['hdr_sz']}pt body={ref['body_sz']}pt"
        )
        break
    return ref


def _add_header_bar(slide, ref: dict | None) -> None:
    """Insert the gradient header bar behind the table header row."""
    import copy
    from pptx.oxml.ns import qn
    from lxml import etree

    spTree = slide.shapes._spTree
    geom = REF_GEOM

    if ref and ref.get("bar_el") is not None:
        bar = copy.deepcopy(ref["bar_el"])
        # unique shape id + canonical name
        cNvPr = bar.find(".//" + qn("p:cNvPr"))
        if cNvPr is not None:
            cNvPr.set("id", str(_next_shape_id(slide)))
            cNvPr.set("name", "Header Gradient")
        # normalize geometry
        xfrm = bar.find(".//" + qn("a:xfrm"))
        if xfrm is not None:
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            if off is not None:
                off.set("x", str(int(geom["left"])))
                off.set("y", str(int(geom["top"])))
            if ext is not None:
                ext.set("cx", str(int(geom["width"])))
                ext.set("cy", str(int(geom["bar_h"])))
        spTree.append(bar)
        return

    # No reference bar in the deck → build the CGI gradient rect from scratch
    stops = "".join(
        f'<a:gs pos="{pos}"><a:srgbClr val="{color}"/></a:gs>'
        for pos, color in CGI_GRADIENT_STOPS
    )
    nsmap = (
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    )
    sp_xml = (
        f"<p:sp {nsmap}>"
        f'<p:nvSpPr><p:cNvPr id="{_next_shape_id(slide)}" name="Header Gradient"/>'
        f"<p:cNvSpPr/><p:nvPr/></p:nvSpPr>"
        f'<p:spPr bwMode="gray">'
        f'<a:xfrm><a:off x="{int(geom["left"])}" y="{int(geom["top"])}"/>'
        f'<a:ext cx="{int(geom["width"])}" cy="{int(geom["bar_h"])}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:gradFill><a:gsLst>{stops}</a:gsLst><a:lin ang="0" scaled="0"/></a:gradFill>'
        f'<a:ln w="9525"><a:noFill/></a:ln><a:effectLst/>'
        f"</p:spPr>"
        f'<p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/><a:lstStyle/><a:p/></p:txBody>'
        f"</p:sp>"
    )
    spTree.append(etree.fromstring(sp_xml))


def _next_shape_id(slide) -> int:
    """Return a shape id not yet used on the slide."""
    used = set()
    for shape in slide.shapes:
        try:
            used.add(shape.shape_id)
        except Exception:
            continue
    return max(used, default=1) + 1


def _update_table_texts_in_place(tbl_shape, table_data: dict) -> bool:
    """
    Update an existing table's cell texts without touching any formatting.
    Returns False when dimensions don't match (caller rebuilds instead).
    """
    headers  = table_data.get("headers") or []
    rows     = table_data.get("rows") or []
    all_rows = table_data.get("all_rows") or ([headers] + rows if headers else rows)
    if not all_rows:
        return False

    tbl = tbl_shape.table
    if len(tbl.rows) != len(all_rows):
        return False
    n_cols = len(tbl.columns)
    if max(len(r) for r in all_rows) != n_cols:
        return False

    for ri, row_data in enumerate(all_rows):
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            new_text = str(row_data[ci]) if ci < len(row_data) else ""
            tf = cell.text_frame
            while len(tf.paragraphs) > 1:
                p = tf.paragraphs[-1]._p
                p.getparent().remove(p)
            para = tf.paragraphs[0]
            runs = para.runs
            if runs:
                runs[0].text = new_text
                for run in runs[1:]:
                    run._r.getparent().remove(run._r)
            elif new_text:
                run = para.add_run()
                run.text = new_text
    logger.debug("[PPTBuilder] Table texts updated in place")
    return True


def _content_weighted_widths(all_rows: list, n_cols: int, total_emu: int) -> list[int]:
    """Distribute table width across columns weighted by their longest content."""
    weights = []
    for ci in range(n_cols):
        longest = 4
        for row in all_rows:
            if ci < len(row):
                longest = max(longest, len(str(row[ci])))
        weights.append(max(5, min(longest, 70)))
    total_w = sum(weights)
    widths = [max(int(total_emu * w / total_w), int(Inches(0.45))) for w in weights]
    # Renormalize: floors/rounding can overshoot — absorb into the widest column
    delta = total_emu - sum(widths)
    widths[widths.index(max(widths))] += delta
    return widths


def _add_reference_table(slide, table_data: dict, ref: dict | None) -> None:
    """
    Add a table in the deck's reference design: uniform geometry, transparent
    style (gradient bar provides the header background), white bold centered
    header text, content-weighted column widths.
    """
    headers  = table_data.get("headers") or []
    rows     = table_data.get("rows") or []
    all_rows = table_data.get("all_rows") or ([headers] + rows if headers else rows)
    if not all_rows:
        return

    num_rows = len(all_rows)
    num_cols = max(len(r) for r in all_rows)
    if num_cols == 0:
        return

    ref = ref or {}
    font_name = ref.get("font", "Arial")
    hdr_sz    = float(ref.get("hdr_sz", 18))
    body_sz   = float(ref.get("body_sz", 11))
    if num_cols >= 7:
        hdr_sz = min(hdr_sz, 16.0)

    geom = REF_GEOM
    slide_h = Inches(7.5)
    avail_h = slide_h - geom["top"] - geom["bottom_margin"]
    hdr_h   = Emu(490000)                  # header row ≈ gradient bar height
    n_data  = max(num_rows - 1, 1)
    data_h  = min(Emu(799000), Emu(int((avail_h - hdr_h) / n_data)))
    tbl_h   = Emu(int(hdr_h) + int(data_h) * n_data)

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols, geom["left"], geom["top"], geom["width"], tbl_h
    )
    try:
        tbl_shape.name = "Table"
    except Exception:
        pass
    tbl = tbl_shape.table

    _set_table_style_id(tbl_shape, ref.get("style_id", TRANSPARENT_TABLE_STYLE_ID))

    widths = _content_weighted_widths(all_rows, num_cols, int(geom["width"]))
    for ci, col in enumerate(tbl.columns):
        col.width = Emu(widths[ci])

    tbl.rows[0].height = Emu(int(hdr_h))
    for ri in range(1, num_rows):
        tbl.rows[ri].height = Emu(int(data_h))

    for ri, row_data in enumerate(all_rows):
        is_header = ri == 0 and bool(headers)
        for ci in range(num_cols):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            for run in list(para.runs):
                run._r.getparent().remove(run._r)
            if is_header:
                para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(row_data[ci]) if ci < len(row_data) else ""

            font = run.font
            font.name = font_name
            font.size = Pt(hdr_sz if is_header else body_sz)
            font.bold = is_header
            if is_header:
                font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    logger.info(
        f"[PPTBuilder] Reference table: {num_rows}r x {num_cols}c "
        f"(style {ref.get('style_id', TRANSPARENT_TABLE_STYLE_ID)})"
    )


def _set_table_style_id(tbl_shape, style_id: str) -> None:
    """Set (replacing any existing) the tableStyleId on a table shape."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree

        tbl_element = tbl_shape._element
        tbl = tbl_element.find(".//" + qn("a:tbl"))
        if tbl is None:
            return
        tblPr = tbl.find(qn("a:tblPr"))
        if tblPr is None:
            tblPr = etree.SubElement(tbl, qn("a:tblPr"))
            tbl.insert(0, tblPr)
        # transparent design: no banding flags, just the style id
        for attr in ("firstRow", "bandRow"):
            if attr in tblPr.attrib:
                del tblPr.attrib[attr]
        existing = tblPr.find(qn("a:tableStyleId"))
        if existing is not None:
            tblPr.remove(existing)
        style_elem = etree.SubElement(tblPr, qn("a:tableStyleId"))
        style_elem.text = style_id
    except Exception as e:
        logger.warning(f"[PPTBuilder] Table style error: {e}")


def _set_placeholder_font_size(ph, size_pt: float) -> None:
    """Set an explicit font size on all runs in a placeholder's text frame."""
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)


def _get_table_position(slide) -> dict:
    """Return position/size of the first TABLE shape on the slide, or defaults."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
    # Default: starts at 0.53" from top (right after compact title strip)
    return {
        "left": Emu(int(Inches(0.53))),
        "top":  Emu(int(Inches(0.79))),
        "width": Emu(int(Inches(12.02))),
        "height": Emu(int(Inches(6.5))),
    }


def _remove_table_shapes(slide) -> None:
    """Remove all TABLE shapes from a slide."""
    to_remove = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE
    ]
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    if to_remove:
        logger.debug(f"[PPTBuilder] Removed {len(to_remove)} table(s) from slide")


def _add_cgi_table(slide, table_data: dict, pos: dict) -> None:
    """
    Add a table to a slide using the CGI table style.
    The CGI style provides the dark header + alternating rows from the theme —
    NO explicit cell fills are needed when using the style.
    Falls back to manual styling if template theme is not available.
    """
    headers  = table_data.get("headers") or []
    rows     = table_data.get("rows") or []
    all_rows = table_data.get("all_rows") or ([headers] + rows if headers else rows)

    if not all_rows:
        return

    num_rows = len(all_rows)
    num_cols = max(len(r) for r in all_rows)
    if num_cols == 0:
        return

    left   = pos["left"]
    top    = pos["top"]
    width  = pos["width"]
    height = pos["height"]

    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    tbl = tbl_shape.table

    # ── Apply CGI table style ─────────────────────────────────────────────────
    # The style GUID makes the theme provide dark header + banded rows.
    _apply_cgi_table_style(tbl_shape, has_header=bool(headers))

    # ── Even column widths ────────────────────────────────────────────────────
    col_width = width // num_cols
    for col in tbl.columns:
        col.width = col_width

    # ── Row height distribution ───────────────────────────────────────────────
    if bool(headers) and num_rows > 1:
        header_row_height = Emu(int(Inches(0.45)))
        data_rows         = num_rows - 1
        data_row_height   = Emu(int((height - header_row_height) / data_rows))
        tbl.rows[0].height = header_row_height
        for ri in range(1, num_rows):
            tbl.rows[ri].height = data_row_height

    # ── Adaptive font sizing ──────────────────────────────────────────────────
    density = num_rows * num_cols
    if num_cols >= 7 or density >= 56:
        table_size  = 10.0
        header_size = 11.0
    elif num_cols >= 6 or density >= 40:
        table_size  = 11.0
        header_size = 12.0
    elif num_rows >= 12 or num_cols >= 5:
        table_size  = 12.0
        header_size = 13.0
    else:
        table_size  = 14.0
        header_size = 14.0

    # ── Fill cells ────────────────────────────────────────────────────────────
    for ri, row_data in enumerate(all_rows):
        is_header = ri == 0 and bool(headers)
        for ci in range(num_cols):
            cell = tbl.cell(ri, ci)
            cell_text = row_data[ci] if ci < len(row_data) else ""

            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            for run in list(para.runs):
                run._r.getparent().remove(run._r)
            run = para.add_run()
            run.text = cell_text

            font = run.font
            font.size = Pt(header_size if is_header else float(table_size))
            font.bold = is_header
            # Do NOT set font.name or font.color.rgb — let theme handle these

    logger.info(f"[PPTBuilder] CGI table: {num_rows}r x {num_cols}c (style applied)")


def _apply_cgi_table_style(tbl_shape, has_header: bool = True) -> None:
    """
    Apply the CGI table style GUID to the table element.
    Removes any existing tableStyleId and inserts the CGI one.
    """
    try:
        from pptx.oxml.ns import qn
        from lxml import etree

        tbl_element = tbl_shape._element
        tblPr = tbl_element.find(".//" + qn("a:tblPr"))
        if tblPr is None:
            tbl = tbl_element.find(".//" + qn("a:tbl"))
            if tbl is None:
                return
            tblPr = etree.SubElement(tbl, qn("a:tblPr"))

        # Set banding attributes
        if has_header:
            tblPr.set("firstRow", "1")
            tblPr.set("bandRow", "1")

        # Remove existing style ID
        existing = tblPr.find(qn("a:tableStyleId"))
        if existing is not None:
            tblPr.remove(existing)

        # Insert CGI style ID
        style_elem = etree.SubElement(tblPr, qn("a:tableStyleId"))
        style_elem.text = CGI_TABLE_STYLE_ID

        logger.debug(f"[PPTBuilder] Applied CGI table style: {CGI_TABLE_STYLE_ID}")
    except Exception as e:
        logger.warning(f"[PPTBuilder] Table style error: {e}")


def _fill_agenda_placeholders(slide, content) -> None:
    """
    Fill the CGI Agenda layout's numbered section placeholders.

    Layout placeholders (index → role):
      ph[0]  = Title ("Agenda")
      ph[18] = Item 1 number  │  ph[16] = Item 1 description
      ph[19] = Item 2 number  │  ph[27] = Item 2 description
      ph[20] = Item 3 number  │  ph[28] = Item 3 description
      ph[26] = Item 4 number  │  ph[29] = Item 4 description
      ph[22] = Item 5 number  │  ph[30] = Item 5 description
      ph[24] = Item 6 number  │  ph[31] = Item 6 description
    """
    # Parse content into a list of section names
    if isinstance(content, list):
        items = [str(item).strip() for item in content if item]
    elif isinstance(content, str):
        items = [
            line.strip()
            for line in content.split("\n")
            if line.strip()
        ]
    else:
        items = []

    # Strip leading numbering ("01. " / "1. " etc.) from items
    import re as _re
    cleaned: list[str] = []
    for item in items:
        item = _re.sub(r"^\d+[.\)]\s*", "", item).strip()
        if item:
            cleaned.append(item)
    items = cleaned

    # Placeholder index pairs: (number_idx, description_idx)
    item_pairs = [(18, 16), (19, 27), (20, 28), (26, 29), (22, 30), (24, 31)]

    ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for i, (num_idx, desc_idx) in enumerate(item_pairs):
        num_ph   = ph_map.get(num_idx)
        desc_ph  = ph_map.get(desc_idx)

        num_text  = f"0{i + 1}" if i < 9 else str(i + 1)
        desc_text = items[i] if i < len(items) else ""

        if num_ph:
            _replace_placeholder_text(num_ph, num_text)
        if desc_ph:
            _replace_placeholder_text(desc_ph, desc_text)


def _update_cover_subtitle(slide, content) -> None:
    """
    Update the subtitle TextBox on the cover slide.
    The CGI cover slide has a manual TextBox (not a placeholder) for the subtitle.
    """
    subtitle_text = ""
    if isinstance(content, list) and content:
        subtitle_text = str(content[0]).strip()
    elif isinstance(content, str):
        subtitle_text = content.strip()

    if not subtitle_text:
        return

    # Find the existing TextBox (shape_type=17 = TEXT_BOX)
    for shape in slide.shapes:
        if shape.shape_type == 17 and hasattr(shape, "text_frame"):
            # This is likely the subtitle textbox (type TEXT_BOX, not a placeholder)
            tf = shape.text_frame
            # Clear extra paragraphs, then set text on the first run so its
            # formatting (size, color, font) is preserved
            while len(tf.paragraphs) > 1:
                p = tf.paragraphs[-1]._p
                p.getparent().remove(p)
            para = tf.paragraphs[0]
            runs = para.runs
            if runs:
                runs[0].text = subtitle_text
                for run in runs[1:]:
                    run._r.getparent().remove(run._r)
            else:
                run = para.add_run()
                run.text = subtitle_text
            logger.debug(f"[PPTBuilder] Cover subtitle updated: {subtitle_text!r}")
            return


def _remove_last_slide(prs: Presentation) -> None:
    """Remove the last slide from a presentation."""
    xml_slides = prs.slides._sldIdLst
    if len(xml_slides) == 0:
        return
    last_sldId = xml_slides[-1]
    rId = last_sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception as e:
            logger.warning(f"[PPTBuilder] drop_rel: {e}")
    xml_slides.remove(last_sldId)


def _is_agenda_slide(data: dict) -> bool:
    """Return True if this slide is the agenda/index slide."""
    title = (data.get("title") or "").lower()
    return "agenda" in title or "index" in title


def _add_new_slide_to_prs(prs, layout_map: dict, slide_data: dict, theme: dict) -> None:
    """Add a new slide to the presentation (used in blank mode or overflow)."""
    layout_name = slide_data.get("layout", "title_content")
    table_data  = slide_data.get("table")

    # For table slides prefer title_only (has only title placeholder)
    if table_data and "title_only" in layout_map:
        layout = layout_map["title_only"]
    else:
        layout = layout_map.get(layout_name) or layout_map.get("title_content") or prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout)
    _apply_background(slide, theme.get("background_color"))
    _populate_slide(slide, slide_data, theme)


# ── Layout map ────────────────────────────────────────────────────────────────

def _build_layout_map(prs: Presentation) -> dict:
    """
    Build a name → layout object mapping from the presentation's slide layouts.
    Handles both standard PowerPoint layout names and CGI-specific custom names.
    """
    layouts = prs.slide_layouts
    mapping = {}
    for layout in layouts:
        name = layout.name.lower()
        # Title / cover slide
        if ("title slide" in name and "content" not in name) or name == "title":
            if "title" not in mapping:
                mapping["title"] = layout
        # Standard content layout
        elif "title and content" in name and not name.startswith("1_"):
            if "title_content" not in mapping:
                mapping["title_content"] = layout
        # CGI "Title only slide" — preferred for table slides (single title placeholder)
        elif "title only slide" in name:
            mapping["title_only"] = layout       # CGI custom
        elif "title only" in name:
            if "title_only" not in mapping:
                mapping["title_only"] = layout   # generic fallback
        # Two-column
        elif "two content" in name or "two_column" in name:
            mapping["two_column"] = layout
        # Blank
        elif "blank" in name:
            mapping["blank"] = layout
        # Section header
        elif "section header" in name:
            mapping["section"] = layout
        # CGI Agenda layout
        elif "agenda" in name:
            mapping["agenda"] = layout

    # Hard fallbacks when no name matched
    if "title" not in mapping and layouts:
        mapping["title"] = layouts[0]
    if "title_content" not in mapping and len(layouts) > 1:
        mapping["title_content"] = layouts[1]
    if "two_column" not in mapping and len(layouts) > 3:
        mapping["two_column"] = layouts[3]
    if "blank" not in mapping and len(layouts) > 6:
        mapping["blank"] = layouts[6]

    logger.debug(f"[PPTBuilder] Layout map keys: {list(mapping.keys())}")
    return mapping


# ── Slide population ──────────────────────────────────────────────────────────

def _populate_slide(slide, data: dict, theme: dict):
    """Fill slide placeholders with title, content/table, and notes."""
    title_text = data.get("title", "")
    content    = data.get("content", "")
    table_data = data.get("table")       # {"headers": [...], "rows": [[...],...]}
    left       = data.get("left_content")
    right      = data.get("right_content")
    notes_text = data.get("notes", "")

    # Resolve theme values
    fonts       = theme.get("fonts") or []
    font_sizes  = theme.get("font_sizes") or []
    text_colors = theme.get("text_colors") or []

    title_font  = fonts[0] if fonts else None
    body_font   = fonts[1] if len(fonts) > 1 else title_font
    title_size  = font_sizes[0] if font_sizes else None
    body_size   = font_sizes[1] if len(font_sizes) > 1 else (font_sizes[0] if font_sizes else None)
    title_color = text_colors[0] if text_colors else None
    body_color  = text_colors[1] if len(text_colors) > 1 else title_color

    # ── Title ─────────────────────────────────────────────────────────────────
    if title_text:
        if slide.shapes.title:
            title_ph = slide.shapes.title

            # For table slides: compact the title placeholder to the top strip
            # so the table can sit close below it (target: 0.25" top, 0.54" height)
            if table_data:
                title_ph.top    = Inches(0.20)
                title_ph.height = Inches(0.54)
                title_ph.left   = Inches(0.5)
                title_ph.width  = Inches(12.33)

            _set_text_frame(
                title_ph.text_frame,
                [title_text],
                font_name=title_font,
                font_size_pt=title_size,
                font_color_hex=title_color,
                bold=True,
                clear=True,
            )
            # Force left-alignment for all content/table slide titles
            if table_data or data.get("layout") == "title_content":
                for para in title_ph.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
        else:
            # Blank layout — add a title text box manually at the top
            _add_title_textbox(
                slide, title_text,
                font_name=title_font,
                font_size_pt=title_size or 28,
                font_color_hex=title_color,
            )

    # ── Table (takes priority over bullet content for data slides) ────────────
    if table_data and not (left or right):
        _add_table_to_slide(slide, table_data, theme)

    # ── Bullet/text content (only when no table and no two-column) ────────────
    elif not (left or right) and not table_data:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                items = content if isinstance(content, list) else ([content] if content else [])
                _set_text_frame(
                    ph.text_frame, items,
                    font_name=body_font, font_size_pt=body_size,
                    font_color_hex=body_color, bold=False, clear=True,
                )
                break

    # ── Two-column layout ─────────────────────────────────────────────────────
    if left or right:
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 1 and left:
                items = left if isinstance(left, list) else ([left] if left else [])
                _set_text_frame(
                    ph.text_frame, items,
                    font_name=body_font, font_size_pt=body_size,
                    font_color_hex=body_color, bold=False, clear=True,
                )
            elif idx == 2 and right:
                items = right if isinstance(right, list) else ([right] if right else [])
                _set_text_frame(
                    ph.text_frame, items,
                    font_name=body_font, font_size_pt=body_size,
                    font_color_hex=body_color, bold=False, clear=True,
                )

    # ── Speaker notes ─────────────────────────────────────────────────────────
    if notes_text:
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
        except Exception as e:
            logger.warning(f"[PPTBuilder] Notes error: {e}")


# ── Title text box (for blank layouts) ───────────────────────────────────────

def _add_title_textbox(slide, title_text: str, *, font_name, font_size_pt, font_color_hex):
    """
    Add a title text box at the top of a slide that has no title placeholder.
    Used for 'blank' layout slides (e.g. table-only slides).
    """
    from pptx.util import Emu
    left   = Inches(0.4)
    top    = Inches(0.15)
    width  = Inches(12.5)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False

    para = tf.paragraphs[0]
    for run in list(para.runs):
        run._r.getparent().remove(run._r)
    run = para.add_run()
    run.text = title_text

    font = run.font
    if font_name:
        font.name = font_name
    font.size = Pt(float(font_size_pt))
    font.bold = True
    rgb = _parse_hex_color(font_color_hex)
    if rgb:
        font.color.rgb = RGBColor(*rgb)
    logger.debug(f"[PPTBuilder] Title textbox added: '{title_text}'")


# ── Table builder ─────────────────────────────────────────────────────────────

def _add_table_to_slide(slide, table_data: dict, theme: dict):
    """
    Add a formatted table to a slide using python-pptx's add_table().
    Positioned below the title with consistent margins.
    Applies theme colors: header row uses title color background, body rows alternate.
    """
    headers  = table_data.get("headers") or []
    rows     = table_data.get("rows") or []
    all_rows = table_data.get("all_rows") or ([headers] + rows if headers else rows)

    if not all_rows:
        return

    num_rows = len(all_rows)
    num_cols = max(len(r) for r in all_rows)
    if num_cols == 0:
        return

    # Fonts / colors from theme
    fonts       = theme.get("fonts") or []
    font_sizes  = theme.get("font_sizes") or []
    text_colors = theme.get("text_colors") or []

    body_font   = fonts[1] if len(fonts) > 1 else (fonts[0] if fonts else "Calibri")
    header_font = fonts[0] if fonts else "Calibri"

    # ── Adaptive font sizing based on table density ───────────────────────────
    # Large tables (many rows OR many columns) need smaller fonts so rows fit
    # within the fixed table height without overflowing / being clipped.
    base_size = float(font_sizes[1] if len(font_sizes) > 1 else (font_sizes[0] if font_sizes else 11))
    density   = num_rows * num_cols          # total cells as proxy for density

    if num_cols >= 7 or density >= 56:       # e.g. 8-col × 8-row = 64 cells
        table_size  = min(base_size, 10.0)
        header_size = 11.0
    elif num_cols >= 6 or density >= 40:     # e.g. 6-col × 7-row = 42 cells
        table_size  = min(base_size, 11.0)
        header_size = 12.0
    elif num_rows >= 12 or num_cols >= 5:    # medium-density
        table_size  = min(base_size, 12.0)
        header_size = 13.0
    else:                                    # small/simple table
        table_size  = min(base_size, 14.0)
        header_size = max(table_size, 14.0)

    body_text_color  = _parse_hex_color(text_colors[1] if len(text_colors) > 1 else (text_colors[0] if text_colors else None))

    # Header row background: use the brand dark color as bg, white text for contrast
    title_color_hex = text_colors[0] if text_colors else "2E4057"
    header_bg = _parse_hex_color(title_color_hex)
    # Auto-contrast: white text on dark bg, dark text on light bg
    if header_bg:
        luminance = (header_bg[0]*299 + header_bg[1]*587 + header_bg[2]*114) / 1000
        header_text_rgb = RGBColor(0xFF, 0xFF, 0xFF) if luminance < 128 else RGBColor(0x1F, 0x1F, 0x1F)
    else:
        header_text_rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Light alternating row bg
    alt_bg = _parse_hex_color("EEF2F7")

    # ── Table position ────────────────────────────────────────────────────────
    left_margin = Inches(0.5)
    tbl_width   = Inches(12.33)          # 13.33" wide slide - 0.5" left - 0.5" right

    if slide.shapes.title:
        # Layout has a real title placeholder — sit 0.05" below it
        title_bottom = slide.shapes.title.top + slide.shapes.title.height
        top_margin = title_bottom + Inches(0.05)
    else:
        # Blank layout — manual title textbox occupies top ~0.85"
        top_margin = Inches(0.85)

    tbl_height = max(Inches(1.0), Inches(7.5) - top_margin - Inches(0.15))

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols, left_margin, top_margin, tbl_width, tbl_height
    )
    tbl = tbl_shape.table

    # Set even column widths
    col_width = tbl_width // num_cols
    for col in tbl.columns:
        col.width = col_width

    # ── Row height distribution ───────────────────────────────────────────────
    # Header row gets a fixed compact height; remaining height split among data rows.
    # This prevents the header from consuming too much vertical space on sparse tables.
    if bool(headers) and num_rows > 1:
        header_row_height = Emu(int(Inches(0.45)))          # ~0.45" for header
        data_rows         = num_rows - 1
        data_row_height   = Emu(int((tbl_height - header_row_height) / data_rows))
        tbl.rows[0].height = header_row_height
        for ri in range(1, num_rows):
            tbl.rows[ri].height = data_row_height

    # Fill cells
    for ri, row_data in enumerate(all_rows):
        is_header = ri == 0 and bool(headers)
        for ci in range(num_cols):
            cell = tbl.cell(ri, ci)
            cell_text = row_data[ci] if ci < len(row_data) else ""

            # Set text
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            # Clear existing runs
            for run in list(para.runs):
                run._r.getparent().remove(run._r)
            run = para.add_run()
            run.text = cell_text

            # Font styling
            font = run.font
            font.name = header_font if is_header else body_font
            font.size = Pt(header_size if is_header else float(table_size))
            font.bold = is_header

            if is_header:
                font.color.rgb = header_text_rgb
            elif body_text_color:
                font.color.rgb = RGBColor(*body_text_color)

            # Cell background
            if is_header and header_bg:
                _set_cell_fill(cell, RGBColor(*header_bg))
            elif not is_header and ri % 2 == 0 and alt_bg:
                _set_cell_fill(cell, RGBColor(*alt_bg))

    logger.info(f"[PPTBuilder] Table added: {num_rows}r x {num_cols}c")


def _set_cell_fill(cell, rgb: RGBColor):
    """Apply a solid fill color to a table cell."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # Remove existing fill
        for child in list(tcPr):
            if child.tag.endswith("}solidFill") or child.tag.endswith("}noFill"):
                tcPr.remove(child)
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", "{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2]))
    except Exception as e:
        logger.warning(f"[PPTBuilder] Cell fill error: {e}")


# ── Text frame helper ─────────────────────────────────────────────────────────

def _set_text_frame(
    tf,
    items: list[str],
    *,
    font_name: str | None,
    font_size_pt: float | None,
    font_color_hex: str | None,
    bold: bool = False,
    clear: bool = True,
):
    """
    Write a list of strings into a text frame, applying font formatting per run.
    Uses add_run() — setting p.text= directly destroys all formatting.
    """
    if clear:
        tf.clear()

    rgb = _parse_hex_color(font_color_hex) if font_color_hex else None

    for i, text in enumerate(items):
        if not text:
            continue
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = 0

        # Remove leftover runs from clear()
        for run in list(para.runs):
            run._r.getparent().remove(run._r)

        run = para.add_run()
        run.text = str(text)

        font = run.font
        if font_name:
            font.name = font_name
        if font_size_pt:
            font.size = Pt(float(font_size_pt))
        font.bold = bold
        if rgb:
            font.color.rgb = RGBColor(*rgb)


# ── Background color ──────────────────────────────────────────────────────────

def _apply_background(slide, color_hex: str | None):
    """Apply solid background color to a slide. Skips if color_hex is None."""
    if not color_hex:
        return
    rgb = _parse_hex_color(color_hex)
    if not rgb:
        return
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)
        logger.debug(f"[PPTBuilder] Background → #{color_hex}")
    except Exception as e:
        logger.warning(f"[PPTBuilder] Background error: {e}")


# ── Color utility ─────────────────────────────────────────────────────────────

def _parse_hex_color(hex_str: str | None) -> tuple[int, int, int] | None:
    """
    Parse '2E4057' or '#2E4057' → (r, g, b) tuple.
    Returns None on any parse error or if input is None/empty.
    """
    if not hex_str:
        return None
    try:
        h = str(hex_str).lstrip("#").strip()
        if len(h) != 6:
            return None
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return None

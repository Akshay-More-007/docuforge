"""
style_parser.py — Extract style and theme metadata from .docx and .pptx files.
Handles fonts, colors, margins, heading styles, slide themes.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# ── Entry point ───────────────────────────────────────────────────────────────

def extract_styles(path: str) -> dict:
    """Auto-detect file type and extract styles."""
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return extract_docx_styles(path)
    elif ext == ".pptx":
        return extract_pptx_styles(path)
    else:
        logger.warning(f"[StyleParser] Unsupported type: {ext}")
        return {}


# ── DOCX ──────────────────────────────────────────────────────────────────────

def extract_docx_styles(path: str) -> dict:
    from docx import Document
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")
    doc = Document(str(p))
    return {
        "file_type": "docx",
        "source_path": str(p),
        "filename": p.name,
        "page_margins": _docx_margins(doc),
        "page_size": _docx_page_size(doc),
        "default_font": _docx_default_font(doc),
        "heading_styles": _docx_heading_styles(doc),
        "normal_style": _docx_paragraph_style(doc, "Normal"),
        "table_style": _docx_table_style(doc),
    }


def _docx_margins(doc) -> dict:
    try:
        s = doc.sections[0]
        return {
            "top_inches": round(s.top_margin.inches, 2),
            "bottom_inches": round(s.bottom_margin.inches, 2),
            "left_inches": round(s.left_margin.inches, 2),
            "right_inches": round(s.right_margin.inches, 2),
        }
    except Exception as e:
        logger.warning(f"[StyleParser] Margins: {e}")
        return {"top_inches": 1.0, "bottom_inches": 1.0, "left_inches": 1.0, "right_inches": 1.0}


def _docx_page_size(doc) -> dict:
    try:
        s = doc.sections[0]
        return {
            "width_inches": round(s.page_width.inches, 2),
            "height_inches": round(s.page_height.inches, 2),
            "orientation": "landscape" if s.page_width > s.page_height else "portrait",
        }
    except Exception as e:
        logger.warning(f"[StyleParser] Page size: {e}")
        return {"width_inches": 8.5, "height_inches": 11.0, "orientation": "portrait"}


def _docx_default_font(doc) -> dict:
    try:
        style = doc.styles["Normal"]
        font = style.font
        return {
            "name": font.name or "Calibri",
            "size_pt": font.size.pt if font.size else 11,
            "bold": bool(font.bold),
            "italic": bool(font.italic),
            "color_rgb": _safe_rgb(font),
        }
    except Exception as e:
        logger.warning(f"[StyleParser] Default font: {e}")
        return {"name": "Calibri", "size_pt": 11, "bold": False, "italic": False, "color_rgb": None}


def _docx_heading_styles(doc) -> dict:
    styles = {}
    for level in range(1, 7):
        name = f"Heading {level}"
        try:
            available = [s.name for s in doc.styles]
            if name not in available:
                continue
            style = doc.styles[name]
            font = style.font
            pf = style.paragraph_format
            styles[name] = {
                "font_name": font.name,
                "size_pt": font.size.pt if font.size else None,
                "bold": font.bold,
                "italic": font.italic,
                "color_rgb": _safe_rgb(font),
                "space_before_pt": pf.space_before.pt if pf.space_before else None,
                "space_after_pt": pf.space_after.pt if pf.space_after else None,
                "keep_with_next": pf.keep_with_next,
            }
        except Exception as e:
            logger.warning(f"[StyleParser] {name}: {e}")
    return styles


def _docx_paragraph_style(doc, name: str) -> dict:
    try:
        style = doc.styles[name]
        pf = style.paragraph_format
        return {
            "line_spacing": str(pf.line_spacing) if pf.line_spacing else "single",
            "space_before_pt": pf.space_before.pt if pf.space_before else 0,
            "space_after_pt": pf.space_after.pt if pf.space_after else 8,
            "alignment": str(pf.alignment) if pf.alignment else "LEFT",
        }
    except Exception as e:
        logger.warning(f"[StyleParser] Para style {name}: {e}")
        return {}


def _docx_table_style(doc):
    try:
        if doc.tables:
            return doc.tables[0].style.name
    except Exception:
        pass
    return None


def _safe_rgb(font):
    try:
        if font.color and font.color.type:
            return str(font.color.rgb)
    except Exception:
        pass
    return None


# ── PPTX ──────────────────────────────────────────────────────────────────────

def extract_pptx_styles(path: str) -> dict:
    from pptx import Presentation
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")
    prs = Presentation(str(p))
    return {
        "file_type": "pptx",
        "source_path": str(p),
        "filename": p.name,
        "slide_width_inches": round(prs.slide_width.inches, 2),
        "slide_height_inches": round(prs.slide_height.inches, 2),
        "slide_count": len(prs.slides),
        "layout_names": _pptx_layout_names(prs),
        "master_fonts": _pptx_master_fonts(prs),
        "slide_themes": _pptx_all_slide_themes(prs),
    }


def extract_pptx_theme(path: str, slide_index: int = 3) -> dict:
    """
    Extract design theme from a specific slide (0-indexed).

    Business PPTs typically store fonts and colors at the slide-master level,
    not on individual runs. Run-level inspection returns "inherit" for everything.
    This function also reads master-level fonts and theme accent colors so the
    LLM always has concrete values to apply.
    """
    from pptx import Presentation
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")
    prs = Presentation(str(p))
    total = len(prs.slides)
    if slide_index >= total:
        slide_index = min(3, total - 1)

    slide = prs.slides[slide_index]
    theme = _pptx_slide_theme(slide, slide_index)

    # Always include master-level fonts (fills in when runs show "inherit")
    master_fonts = _pptx_master_fonts(prs, str(p))
    theme["master_fonts"] = master_fonts

    # Extract theme accent colors from XML (fills in when no direct RGB on runs)
    theme["theme_accent_colors"] = _pptx_theme_accent_colors(prs, str(p))

    # Resolve concrete font names: prefer run-level, fall back to master
    if not theme["fonts"]:
        resolved = []
        if master_fonts.get("title"):
            resolved.append(master_fonts["title"])
        if master_fonts.get("body"):
            resolved.append(master_fonts["body"])
        if resolved:
            theme["fonts"] = resolved
            logger.info(f"[StyleParser] Fonts resolved from master: {resolved}")

    theme["total_slides"] = total
    theme["note"] = (
        f"Theme extracted from slide {slide_index + 1} of {total}. "
        f"Master fonts: title={master_fonts.get('title')}, body={master_fonts.get('body')}. "
        f"Apply fonts/colors from master_fonts and theme_accent_colors if fonts/text_colors lists are empty."
    )
    return theme


def _pptx_theme_accent_colors(prs, pptx_path: str | None = None) -> dict:
    """
    Extract the theme accent colors (dk1, lt1, accent1-6) from ppt/theme/theme1.xml.

    The clrScheme lives in the theme XML file inside the PPTX zip, NOT in
    master.element (which is a different XML tree). We read it via zipfile,
    with master.element as a fallback. Returns hex strings without '#'.
    """
    colors: dict = {}
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Primary: read directly from theme XML via zipfile
    if pptx_path:
        try:
            import zipfile
            from lxml import etree
            with zipfile.ZipFile(pptx_path) as z:
                theme_files = sorted(
                    [n for n in z.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
                )
                if theme_files:
                    xml = z.read(theme_files[0])
                    root = etree.fromstring(xml)
                    clrScheme = root.find(f".//{{{ns}}}clrScheme")
                    if clrScheme is not None:
                        for child in clrScheme:
                            tag = child.tag.split("}")[-1]  # dk1, lt1, accent1 ...
                            srgb = child.find(f"{{{ns}}}srgbClr")
                            sys_clr = child.find(f"{{{ns}}}sysClr")
                            if srgb is not None:
                                colors[tag] = srgb.get("val", "").upper()
                            elif sys_clr is not None:
                                colors[tag] = sys_clr.get("lastClr", "").upper()
            if colors:
                return colors
        except Exception as e:
            logger.warning(f"[StyleParser] Theme accent colors (zip): {e}")

    # Fallback: search master element (works for some PPTXs)
    try:
        master = prs.slide_master
        clrScheme = master.element.find(f".//{{{ns}}}clrScheme")
        if clrScheme is not None:
            for child in clrScheme:
                tag = child.tag.split("}")[-1]
                srgb = child.find(f"{{{ns}}}srgbClr")
                sys_clr = child.find(f"{{{ns}}}sysClr")
                if srgb is not None:
                    colors[tag] = srgb.get("val", "").upper()
                elif sys_clr is not None:
                    colors[tag] = sys_clr.get("lastClr", "").upper()
    except Exception as e:
        logger.warning(f"[StyleParser] Theme accent colors (master): {e}")

    return colors


def _pptx_layout_names(prs) -> list:
    try:
        return [layout.name for layout in prs.slide_layouts]
    except Exception:
        return []


def _pptx_master_fonts(prs, pptx_path: str | None = None) -> dict:
    """
    Extract master font names from the theme XML inside the PPTX zip.

    python-pptx's master.element doesn't contain the fontScheme directly —
    it lives in ppt/theme/theme1.xml. We read it via zipfile to get the
    actual typeface names (e.g. 'Aptos Display', 'Aptos', 'Calibri').
    """
    fonts = {"title": None, "body": None}

    # Try direct zip read first (most reliable)
    if pptx_path:
        try:
            import zipfile
            from lxml import etree
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            with zipfile.ZipFile(pptx_path) as z:
                theme_files = sorted(
                    [n for n in z.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
                )
                if theme_files:
                    xml = z.read(theme_files[0])
                    root = etree.fromstring(xml)
                    fe = root.find(f".//{{{ns}}}fontScheme")
                    if fe is not None:
                        major = fe.find(f"{{{ns}}}majorFont/{{{ns}}}latin")
                        minor = fe.find(f"{{{ns}}}minorFont/{{{ns}}}latin")
                        if major is not None:
                            fonts["title"] = major.get("typeface")
                        if minor is not None:
                            fonts["body"] = minor.get("typeface")
            if fonts["title"] or fonts["body"]:
                return fonts
        except Exception as e:
            logger.warning(f"[StyleParser] Direct theme zip read failed: {e}")

    # Fallback: search master element (works if fonts are inlined there)
    try:
        master = prs.slide_master
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        fe = master.element.find(f".//{{{ns}}}fontScheme")
        if fe is not None:
            major = fe.find(f"{{{ns}}}majorFont/{{{ns}}}latin")
            minor = fe.find(f"{{{ns}}}minorFont/{{{ns}}}latin")
            if major is not None:
                fonts["title"] = major.get("typeface")
            if minor is not None:
                fonts["body"] = minor.get("typeface")
    except Exception as e:
        logger.warning(f"[StyleParser] Master fonts fallback: {e}")

    return fonts


def _pptx_all_slide_themes(prs) -> list:
    themes = []
    for i, slide in enumerate(prs.slides):
        try:
            themes.append(_pptx_slide_theme(slide, i))
        except Exception as e:
            logger.warning(f"[StyleParser] Slide {i+1}: {e}")
            themes.append({"slide_index": i, "error": str(e)})
    return themes


def _pptx_slide_theme(slide, index: int) -> dict:
    theme = {
        "slide_index": index,
        "slide_number": index + 1,
        "background_color": _pptx_background_color(slide),
        "fonts": [],
        "font_sizes": [],
        "text_colors": [],
        "shapes": [],
    }
    for shape in slide.shapes:
        info = _pptx_shape_info(shape)
        if info:
            theme["shapes"].append(info)
            for font in info.get("fonts", []):
                if font and font not in theme["fonts"]:
                    theme["fonts"].append(font)
            for size in info.get("font_sizes", []):
                if size and size not in theme["font_sizes"]:
                    theme["font_sizes"].append(size)
            for color in info.get("text_colors", []):
                if color and color not in theme["text_colors"]:
                    theme["text_colors"].append(color)
    return theme


def _pptx_background_color(slide):
    try:
        fill = slide.background.fill
        if fill.type is not None:
            try:
                return str(fill.fore_color.rgb)
            except Exception:
                pass
    except Exception:
        pass
    return None


def _pptx_shape_info(shape):
    try:
        if not hasattr(shape, "text_frame"):
            return None
        info = {
            "shape_name": shape.name,
            "fonts": [], "font_sizes": [], "text_colors": [], "bold": [], "alignment": [],
        }
        for para in shape.text_frame.paragraphs:
            try:
                if para.alignment:
                    align = str(para.alignment)
                    if align not in info["alignment"]:
                        info["alignment"].append(align)
            except Exception:
                pass
            for run in para.runs:
                try:
                    font = run.font
                    if font.name and font.name not in info["fonts"]:
                        info["fonts"].append(font.name)
                    if font.size:
                        size_pt = round(font.size.pt, 1)
                        if size_pt not in info["font_sizes"]:
                            info["font_sizes"].append(size_pt)
                    if font.bold and font.bold not in info["bold"]:
                        info["bold"].append(font.bold)
                    try:
                        if font.color and font.color.type:
                            color = str(font.color.rgb)
                            if color not in info["text_colors"]:
                                info["text_colors"].append(color)
                    except Exception:
                        pass
                except Exception:
                    continue
        return info if (info["fonts"] or info["font_sizes"] or info["text_colors"]) else None
    except Exception:
        return None
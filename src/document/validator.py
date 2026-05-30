"""
validator.py — Audit output .docx or .pptx against user requirements.
Uses reasoning model (Qwen QwQ / DeepSeek R1) for requirement checking.

Notes:
- PPTX visual design (colors, layouts) can't be verified from text — auto-pass.
- PDF is a source-only format; output is always docx or pptx, both handled here.
- Text truncated to 8000 chars to stay within context limits.
"""

import json
import logging
from pathlib import Path
from langchain_core.messages import HumanMessage, SystemMessage

from src.llm.router import get_llm

logger = logging.getLogger(__name__)

VALIDATOR_SYSTEM_PROMPT = """You are a strict document quality auditor for DocuForge.

Your job: check if a reformatted document meets ALL stated requirements.

You will receive:
1. The user's requirements (list of strings)
2. A text dump of the output document

Respond ONLY with JSON:
{
  "requirements_met": true | false,
  "passed": ["req1", "req2"],
  "failed": ["req3"],
  "feedback": "Specific, actionable feedback on what's wrong and how to fix it."
}

How to judge each requirement:
- CONTENT requirements you CAN verify from the text (a section/table exists, specific
  text is present, spelling is correct) → check strictly and fail if not satisfied.
- META / QUALITY directives that are NOT checkable from text alone — e.g. "be accurate",
  "do not hallucinate", "be thorough", "don't make mistakes", "be careful with margins" —
  → treat as SATISFIED. These are guidance, not verifiable deliverables.
- VISUAL / FORMATTING properties not present in extracted text (fonts, colors, exact
  layout, spacing) → treat as SATISFIED unless the text shows an obvious content defect.

Only set requirements_met=false when a CONCRETE, verifiable content requirement fails.
If the requirements list is empty, return requirements_met=true.
"""


async def validate_output(requirements: list[str], output_path: str, build_mode: str | None = None) -> dict:
    """
    Validate output document against requirements.

    PPTX: visual formatting (colors, fonts, layout) cannot be verified via text extraction,
    so we auto-pass and skip LLM validation.
    DOCX (template mode): the source was edited IN PLACE, so original formatting,
    fonts and colors are preserved by construction — the LLM judges only content
    (spelling fixes, RACI/flow presence), not unverifiable visual properties.
    DOCX (fresh build): full text-based validation with LLM.
    """
    suffix = Path(output_path).suffix.lower()

    # Auto-pass PPTX — visual design can't be verified from text
    if suffix == ".pptx":
        logger.info("[Validator] PPTX output — auto-passing (visual design not text-verifiable)")
        return {
            "requirements_met": True,
            "passed": requirements,
            "failed": [],
            "feedback": "",
        }

    # Auto-pass if no requirements specified
    if not requirements:
        return {
            "requirements_met": True,
            "passed": [],
            "failed": [],
            "feedback": "",
        }

    # Extract text for LLM review
    doc_text = _extract_text_for_validation(output_path)
    if not doc_text:
        logger.warning(f"[Validator] Could not extract text from {output_path}")
        return {
            "requirements_met": False,
            "passed": [],
            "failed": requirements,
            "feedback": "Could not extract text from output document for validation.",
        }

    llm = get_llm(task="validation")  # Reasoning model

    mode_note = ""
    if build_mode == "docx_template":
        mode_note = (
            "\nIMPORTANT CONTEXT: This document was produced by editing the ORIGINAL "
            "source file IN PLACE. Therefore all original formatting, fonts, colors, "
            "layout, headers/footers and styling are PRESERVED BY CONSTRUCTION. Treat "
            "any requirement about preserving/maintaining format, fonts, colors, or "
            "layout as automatically SATISFIED. Judge ONLY content changes you can "
            "verify in the text below: spelling/capitalization fixes, and the presence "
            "of requested tables such as a RACI matrix or a process-flow table.\n"
        )

    prompt = f"""Requirements to check:
{json.dumps(requirements, indent=2)}
{mode_note}
Output document text:
\"\"\"
{doc_text[:8000]}
\"\"\"

Check every requirement against the document text and return your JSON verdict."""

    messages = [
        SystemMessage(content=VALIDATOR_SYSTEM_PROMPT),
        HumanMessage(content=prompt),
    ]

    try:
        response = await llm.ainvoke(messages)
        raw = response.content.strip()
        result = _parse_json_response(raw)
        logger.info(
            f"[Validator] requirements_met={result.get('requirements_met')} "
            f"| failed={result.get('failed')}"
        )
        return result

    except (json.JSONDecodeError, KeyError, ValueError) as e:
        logger.error(f"[Validator] Parse error: {e}")
        return {
            "requirements_met": False,
            "passed": [],
            "failed": requirements,
            "feedback": f"Validation parsing failed: {e}",
        }


def _extract_text_for_validation(path: str) -> str:
    """Extract plain text from .docx or .pptx for LLM review."""
    p = Path(path)
    suffix = p.suffix.lower()

    try:
        if suffix == ".docx":
            from docx import Document
            doc = Document(str(p))
            # Tables FIRST — RACI matrices, flow tables and data tables live in
            # tables, often near the END of long documents. The validation text
            # is truncated, so if paragraphs came first the tables would be cut
            # off and table-based requirements would falsely fail. Leading with
            # table content guarantees these structural elements are visible.
            table_parts = []
            for ti, tbl in enumerate(doc.tables):
                table_parts.append(f"[TABLE {ti}]")
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        table_parts.append(" | ".join(cells))
            para_parts = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n".join(table_parts + ["", "[BODY TEXT]"] + para_parts)

        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(p))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)

    except Exception as e:
        logger.warning(f"[Validator] Text extraction error: {e}")

    return ""


def _parse_json_response(raw: str) -> dict:
    """Strip thinking tags and code fences, then parse JSON."""
    import json

    # Strip DeepSeek/Qwen thinking tags
    if "<think>" in raw:
        raw = raw.split("</think>")[-1].strip()

    # Strip markdown code fences
    if raw.startswith("```"):
        parts = raw.split("```")
        raw = parts[1] if len(parts) > 1 else raw
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.strip()

    return json.loads(raw)

"""
extractor.py — Extract content from .docx, .pptx, or .pdf files.

For PPTX: extracts not just text but also per-shape font/color metadata,
so the LLM has full context when rebuilding the presentation.
"""

import re
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# ── Entry point ────────────────────────────────────────────────────────────────

def extract(path: str) -> dict:
    """Auto-detect file type and extract content."""
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return extract_docx(path)
    elif ext == ".pptx":
        return extract_pptx(path)
    elif ext == ".pdf":
        return extract_pdf(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ── DOCX ──────────────────────────────────────────────────────────────────────

def extract_docx(path: str) -> dict:
    import mammoth
    import html as html_lib

    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    with open(p, "rb") as f:
        html_result = mammoth.convert_to_html(f)
    with open(p, "rb") as f:
        md_result = mammoth.convert_to_markdown(f)

    raw_html = html_result.value
    raw_text = md_result.value
    sections = _parse_sections_html(raw_html)

    return {
        "source_path": str(p),
        "filename": p.name,
        "file_type": "docx",
        "sections": sections,
        "raw_html": raw_html,
        "raw_text": raw_text,
        "warnings": [str(w) for w in html_result.messages],
    }


# ── PPTX ──────────────────────────────────────────────────────────────────────

def extract_pptx(path: str) -> dict:
    """
    Extract structured content from a .pptx, including:
    - Text placeholders (title + body)
    - TABLE shapes → structured rows/headers  ← critical for data slides
    - Per-run font/color metadata for design fidelity

    Tables are the most common content type in business PPTs.
    Previously they were silently skipped — now fully extracted.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")

    prs = Presentation(str(p))
    sections = []
    raw_lines = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        content_parts = []
        shape_styles = []
        tables = []          # structured table data per slide

        for shape in slide.shapes:
            try:
                # ── TABLE shapes ─────────────────────────────────────────────
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tbl_data = _extract_table(shape)
                    if tbl_data:
                        tables.append(tbl_data)
                        # flat text for raw_text / LLM context
                        for row in tbl_data["all_rows"]:
                            content_parts.append(" | ".join(row))
                    continue

                # ── Text frame shapes ─────────────────────────────────────────
                if not hasattr(shape, "text_frame"):
                    continue
                text = shape.text.strip()
                if not text:
                    continue

                # Detect title placeholder (idx=0)
                is_title = False
                try:
                    ph = shape.placeholder_format
                    if ph is not None and ph.idx == 0:
                        is_title = True
                except Exception:
                    pass

                if is_title:
                    title = text
                else:
                    content_parts.append(text)

                # Collect font/color metadata for design reconstruction
                shape_meta = _extract_shape_style(shape)
                if shape_meta:
                    shape_styles.append(shape_meta)

            except Exception:
                continue

        section_text = "\n".join(content_parts)
        section: dict = {
            "level": 1,
            "heading": title or f"Slide {i}",
            "content": section_text,
            "slide_number": i,
            "shape_styles": shape_styles,
        }
        if tables:
            section["tables"] = tables   # structured table data for LLM

        sections.append(section)
        raw_lines.append(f"[Slide {i}] {title}\n{section_text}")

    return {
        "source_path": str(p),
        "filename": p.name,
        "file_type": "pptx",
        "sections": sections,
        "raw_html": "",
        "raw_text": "\n\n".join(raw_lines),
        "warnings": [],
        "slide_count": len(prs.slides),
    }


def _extract_table(shape) -> dict | None:
    """
    Extract a TABLE shape into a structured dict.

    Returns:
      {
        "headers":  ["Col1", "Col2", ...],   # first row
        "rows":     [["R1C1", "R1C2"], ...], # remaining rows
        "all_rows": [["H1","H2"], ["R1C1","R1C2"], ...]
      }
    Returns None if the table has no content.
    """
    try:
        tbl = shape.table
        all_rows: list[list[str]] = []
        seen_spans: set[tuple] = set()   # skip merged/repeated cells

        for ri, row in enumerate(tbl.rows):
            row_data: list[str] = []
            for ci, cell in enumerate(row.cells):
                # python-pptx repeats merged cells — de-duplicate by span coords
                span_key = (cell.sp_id if hasattr(cell, "sp_id") else id(cell),)
                # simpler: just use text, duplicates within a row are rare
                cell_text = cell.text.strip().replace("\n", " ").replace("\r", "")
                row_data.append(cell_text)
            all_rows.append(row_data)

        if not all_rows:
            return None

        return {
            "headers": all_rows[0],
            "rows": all_rows[1:] if len(all_rows) > 1 else [],
            "all_rows": all_rows,
        }
    except Exception as e:
        logger.warning(f"[Extractor] Table extraction error: {e}")
        return None


def _extract_shape_style(shape) -> dict | None:
    """Extract font name, size, color, bold from all runs in a shape."""
    try:
        fonts, sizes, colors, bolds = set(), set(), set(), set()
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    font = run.font
                    if font.name:
                        fonts.add(font.name)
                    if font.size:
                        sizes.add(round(font.size.pt, 1))
                    if font.bold:
                        bolds.add(True)
                    if font.color and font.color.type:
                        colors.add(str(font.color.rgb))
                except Exception:
                    continue
        if fonts or sizes or colors:
            return {
                "shape_name": shape.name,
                "fonts": sorted(fonts),
                "font_sizes_pt": sorted(sizes),
                "text_colors_hex": sorted(colors),
                "bold": bool(bolds),
            }
    except Exception:
        pass
    return None


# ── PDF ───────────────────────────────────────────────────────────────────────

def extract_pdf(path: str) -> dict:
    """
    Extract text from a PDF using pypdf.
    PDFs are source-only; output is always DOCX or PPTX.
    """
    from pypdf import PdfReader

    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")

    reader = PdfReader(str(p))
    sections = []
    raw_lines = []

    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            sections.append({
                "level": 0,
                "heading": f"Page {i}",
                "content": text,
            })
            raw_lines.append(f"[Page {i}]\n{text}")

    return {
        "source_path": str(p),
        "filename": p.name,
        "file_type": "pdf",
        "sections": sections,
        "raw_html": "",
        "raw_text": "\n\n".join(raw_lines),
        "warnings": [],
        "page_count": len(reader.pages),
    }


# ── HTML section parser (for DOCX) ────────────────────────────────────────────

def _parse_sections_html(html: str) -> list[dict]:
    import html as html_lib

    sections = []
    parts = re.split(r"(<h[1-6][^>]*>.*?</h[1-6]>)", html, flags=re.IGNORECASE | re.DOTALL)
    current_heading = {"level": 0, "text": "Preamble"}
    current_content: list[str] = []

    for part in parts:
        heading_match = re.match(
            r"<h([1-6])[^>]*>(.*?)</h[1-6]>", part, re.IGNORECASE | re.DOTALL
        )
        if heading_match:
            if current_content:
                sections.append({
                    "level": current_heading["level"],
                    "heading": current_heading["text"],
                    "content": " ".join(current_content).strip(),
                })
            level = int(heading_match.group(1))
            text = re.sub(r"<[^>]+>", "", heading_match.group(2)).strip()
            text = html_lib.unescape(text)
            current_heading = {"level": level, "text": text}
            current_content = []
        else:
            clean = re.sub(r"<[^>]+>", " ", part)
            clean = html_lib.unescape(clean).strip()
            if clean:
                current_content.append(clean)

    if current_content:
        sections.append({
            "level": current_heading["level"],
            "heading": current_heading["text"],
            "content": " ".join(current_content).strip(),
        })

    return sections

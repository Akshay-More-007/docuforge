"""Test using source PPTX as template, clearing slides, then adding new ones."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SOURCE = r"C:\Users\Batman\Downloads\22 CGI_Vendor_Final-2 (002).pptx"
OUT    = r"C:\Users\Batman\Downloads\template_test.pptx"

def remove_all_slides(prs):
    """Remove all slides from presentation, keeping layouts and theme."""
    xml_slides = prs.slides._sldIdLst
    presentation_part = prs.part
    n = len(list(xml_slides))
    for sldId in list(xml_slides):
        rId = sldId.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        if rId:
            try:
                presentation_part.drop_rel(rId)
            except Exception as e:
                print(f"  drop_rel warning: {e}")
        xml_slides.remove(sldId)
    print(f"Removed {n} slides. Remaining: {len(prs.slides)}")


def build_layout_map(prs):
    layouts = prs.slide_layouts
    mapping = {}
    for i, layout in enumerate(layouts):
        name = layout.name.lower()
        print(f"  [{i}] {layout.name!r} -> ", end="")
        if "title slide" in name and "content" not in name:
            mapping["title"] = layout
            print("-> title")
        elif "title and content" in name and "1_" not in name:
            if "title_content" not in mapping:
                mapping["title_content"] = layout
            print("-> title_content")
        elif "two content" in name:
            mapping["two_column"] = layout
            print("-> two_column")
        elif "blank" in name:
            mapping["blank"] = layout
            print("-> blank")
        elif "section header" in name:
            mapping["section"] = layout
            print("-> section")
        elif "title only slide" in name:
            mapping["title_only"] = layout
            print("-> title_only (CGI)")
        elif "title only" in name:
            if "title_only" not in mapping:
                mapping["title_only"] = layout
            print("-> title_only (generic)")
        elif "agenda" in name:
            mapping["agenda"] = layout
            print("-> agenda")
        else:
            print("-> (unmapped)")

    # Fallbacks
    if "title" not in mapping and layouts:
        mapping["title"] = layouts[0]
    if "title_content" not in mapping and len(layouts) > 1:
        mapping["title_content"] = layouts[1]
    if "blank" not in mapping and len(layouts) > 6:
        mapping["blank"] = layouts[6]

    return mapping


# Open source as template
prs = Presentation(SOURCE)
print(f"Opened: {SOURCE}")
print(f"  Slides: {len(prs.slides)}, Layouts: {len(prs.slide_layouts)}")
print(f"  Size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f}")

# Clear slides
remove_all_slides(prs)

# Build layout map
print("\n=== Layout map ===")
lmap = build_layout_map(prs)
print(f"\nMapped: {list(lmap.keys())}")

# Add a test cover slide
cover_layout = lmap.get("title") or prs.slide_layouts[1]
slide = prs.slides.add_slide(cover_layout)
print(f"\nAdded cover slide with layout: {cover_layout.name!r}")
print(f"  Placeholders: {[ph.name for ph in slide.placeholders]}")

# Fill title
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "TA Vendor Empanelment"
        tf = ph.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x0E, 0x28, 0x41)
        break

# Add agenda slide
agenda_layout = lmap.get("agenda")
if agenda_layout:
    aslide = prs.slides.add_slide(agenda_layout)
    print(f"\nAdded agenda slide with layout: {agenda_layout.name!r}")
    items = ["Active Recruitment Vendors", "Active Non-Recruitment Vendors", "Vendor Pipeline – RFP"]

    # ph[0]=title, ph[18]=num1, ph[16]=desc1, ph[19]=num2, ph[27]=desc2, ph[20]=num3, ph[28]=desc3
    item_pairs = [(18, 16), (19, 27), (20, 28)]

    for ph in aslide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "Agenda"
        else:
            for i, (num_idx, desc_idx) in enumerate(item_pairs):
                if idx == num_idx:
                    ph.text = f"0{i+1}"
                    break
                elif idx == desc_idx:
                    if i < len(items):
                        ph.text = items[i]
                    break
    print("  Agenda items set")

# Add a table slide
title_only_layout = lmap.get("title_only")
if title_only_layout:
    tslide = prs.slides.add_slide(title_only_layout)
    print(f"\nAdded table slide with layout: {title_only_layout.name!r}")
    print(f"  Placeholders: {[ph.name for ph in tslide.placeholders]}")

    # Set title
    for ph in tslide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Recruitment Vendors – Non RFP Empanelment Status"
            ph.top = Emu(int(Inches(0.25)))
            ph.height = Emu(int(Inches(0.54)))
            ph.left = Emu(int(Inches(0.53)))
            ph.width = Emu(int(Inches(12.02)))
            tf = ph.text_frame
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x0E, 0x28, 0x41)
            break

    # Add table
    left = Emu(int(Inches(0.53)))
    top = Emu(int(Inches(0.79)))
    width = Emu(int(Inches(12.02)))
    height = Emu(int(Inches(6.5)))

    tbl_shape = tslide.shapes.add_table(4, 4, left, top, width, height)
    tbl = tbl_shape.table
    headers = ["Sl No", "Vendor Name", "Category", "Status"]
    rows_data = [
        ["1", "Test Vendor A", "Permanent", "WIP"],
        ["2", "Test Vendor B", "Contract", "Completed"],
        ["3", "Test Vendor C", "Permanent", "WIP"],
    ]
    all_data = [headers] + rows_data
    for ri, row in enumerate(all_data):
        is_header = ri == 0
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            if is_header:
                from pptx.oxml.ns import qn
                from lxml import etree
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", "0E2841")
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
    print("  Table added")

# Save
prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"  Total slides: {len(prs.slides)}")

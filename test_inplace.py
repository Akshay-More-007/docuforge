"""Test the in-place template modification approach."""
import sys
sys.path.insert(0, ".")
from src.document.ppt_builder import build_pptx

SOURCE = r"C:\Users\Batman\Downloads\22 CGI_Vendor_Final-2 (002).pptx"
OUT    = r"C:\Users\Batman\Downloads\inplace_test.pptx"

# Simulate what the LLM would return for this deck
slides = [
    {
        "layout": "title",
        "title": "TA Vendor Empanelment",
        "content": "Active Vendor Status for FY 2026",
        "slide_number": 1,
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [36, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
    {
        "layout": "title_content",
        "title": "Agenda",
        "content": ["Active Recruitment Vendors", "Active Non-Recruitment Vendors", "Vendor Pipeline – RFP"],
        "slide_number": 2,
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [36, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
    {
        "layout": "title_content",
        "title": "Recruitment Vendors – Non RFP Empanelment Status",
        "slide_number": 3,
        "table": {
            "headers": ["Sl No", "Vendor Name", "Category", "Source", "TA Evaluation", "Procurement", "DPSC", "Status"],
            "rows": [
                ["1", "Resmera Solutions", "Permanent/Contract", "E&U", "WIP", "WIP", "Completed", "Assessment pending from TA."],
                ["2", "Precision Staffers", "Permanent", "PNC/BFS", "Completed", "WIP", "Completed", "TA Evaluation Score-Positive."],
                ["3", "Innovation Roots", "Contract", "I&H", "Not Applicable", "WIP", "WIP", "DPSC call to be scheduled 22nd May"],
                ["4", "EIT", "Permanent", "PNC/BFS", "Completed", "WIP", "Completed", "TA Evaluation Score-Positive."],
                ["5", "INMORPHIS", "Contract", "CME", "WIP", "WIP", "Completed", "Assessment pending from TA."],
                ["6", "InOpTra", "Permanent", "MIG", "WIP", "WIP", "WIP", "Initial discussion pending with Procurement."],
                ["7", "The Merakis Ventures", "Contract", "TA", "WIP", "WIP", "Completed", "TA initial discussion done 21st May"],
            ]
        },
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [28, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
    {
        "layout": "title_content",
        "title": "Active Non-Recruitment Vendors",
        "slide_number": 4,
        "table": {
            "headers": ["Sl No", "Vendor Name", "Category", "Procurement", "DPSC", "Status Note"],
            "rows": [
                ["1", "Callify", "Interview Scheduling", "WIP", "Completed", "Call scheduled on 25th May with Pradeep to discuss certain points"],
                ["2", "Zensible", "TA Tools & Automation", "WIP", "WIP", "Call scheduled on 25th May with Pradeep to discuss certain points"],
                ["3", "Curatal", "Sourcing Channel", "WIP", "Completed", "Rate card finalised 2. MSA in progress. 3. Vendor ID raised."],
                ["4", "Zyamam", "API Integration", "Not Applicable", "Completed", "Security Review Pending"],
                ["5", "Flocareer", "Interview Assessment", "Not Applicable", "Not Applicable", "Approval in place need to check with legal for the next step."],
                ["6", "Naukri", "Job Board", "Not Applicable", "Not Applicable", "NDA shared with Naukri team to review"],
            ]
        },
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [28, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
    {
        "layout": "title_content",
        "title": "Recruitment Vendor – RFP",
        "slide_number": 5,
        "table": {
            "headers": ["SL NO", "Vendor Name", "Category", "Request Source", "Incumbent/New"],
            "rows": [
                ["1", "Srinav Info Systems Private Limited", "Permanent/Contract/Manage services", "TA", "Incumbent Vendors"],
                ["2", "Han Digital Solution Private Limited", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["3", "Live Connections", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["4", "PELATIS RIGAS", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["5", "Golden Opportunities", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["6", "Trigent Software Private Limited", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["7", "Sacha Global", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["8", "Job World India Private Limited", "Permanent", "TA", "Incumbent Vendors"],
                ["9", "I square Soft (Xevyte)", "Permanent/Contract", "TA", "Incumbent Vendors"],
                ["10", "Black and White Business Solutions Pvt Ltd", "Permanent", "TA", "Incumbent Vendors"],
                ["11", "Scaleneworks People Solutions LLP", "Permanent/Manage services", "TA", "Incumbent Vendors"],
                ["12", "Orcapod Consulting Services Private Limited", "Permanent/Contract", "TA", "Incumbent Vendors"],
            ]
        },
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [28, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
    {
        "layout": "title_content",
        "title": "Recruitment Vendor – RFP",
        "slide_number": 6,
        "table": {
            "headers": ["SL NO", "Vendor Name", "Category", "Request Source", "Incumbent/New"],
            "rows": [
                ["13", "Talent Ola Inc.", "Permanent staffing", "TA", "New"],
                ["14", "Careernet Technologies Private Limited", "Permanent staffing", "TA", "New"],
                ["15", "Pacific Consulting", "Permanent staffing", "TA", "New"],
                ["16", "Joan Technologies Private Limited", "Permanent staffing", "TA", "New"],
                ["17", "Nilasu Consulting Services Private Limited", "Permanent staffing", "TA", "New"],
                ["18", "Avtar- Flexi Careers India Private Limited", "Diversity Hiring & HTD", "TA", "New"],
                ["19", "Trailblaze CareerPaths Global Pvt.Ltd.", "Permanent staffing", "TA", "New"],
                ["20", "PEOPLELOGIC BUSINESS SOLUTIONS PRIVATE LIMITED", "Managed services", "TA", "New"],
                ["21", "Nityo Infotech Services Private Limited", "Permanent staffing", "Procurement", "New"],
            ]
        },
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [28, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
    {
        "layout": "title_content",
        "title": "Recruitment Vendor – RFP",
        "slide_number": 7,
        "table": {
            "headers": ["SL NO", "Vendor Name", "Category", "Request Source", "Incumbent/New"],
            "rows": [
                ["22", "Prosperix Innovation Private Limited", "Aggregator", "TA", "New"],
                ["23", "CBREX Technologies Private Limited", "Aggregator", "TA", "New"],
                ["24", "Workflexi", "GIG workforce", "TA", "New"],
                ["25", "2COMS", "GIG workforce", "TA", "New"],
                ["26", "Michael Page", "GIG workforce", "TA", "New"],
                ["27", "Quess Corp", "GIG workforce", "TA", "New"],
                ["28", "SRS Business Solutions India Private Limited", "Permanent staffing", "Procurement", "New"],
                ["29", "Openteq Consultants Private Limited", "Permanent staffing", "Procurement", "New"],
            ]
        },
        "theme": {
            "fonts": ["Aptos Display", "Aptos"],
            "font_sizes": [28, 18],
            "text_colors": ["0E2841", "000000"],
            "background_color": None
        }
    },
]

print(f"Building PPTX from source template: {SOURCE}")
result = build_pptx(slides, OUT, template_path=SOURCE)
print(f"Saved: {result}")

# Verify
from pptx import Presentation
p = Presentation(result)
print(f"Output slides: {len(p.slides)}")
for i, slide in enumerate(p.slides, 1):
    shapes = [(s.name, s.shape_type) for s in slide.shapes]
    print(f"  Slide {i}: {len(shapes)} shapes")
    for name, stype in shapes:
        print(f"    [{stype}] {name!r}")

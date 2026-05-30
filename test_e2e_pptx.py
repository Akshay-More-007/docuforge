"""
End-to-end test: run DocuForge on the real source PPTX.
Mimics exactly what the Streamlit app does.
"""
import asyncio
import sys
import logging
import os
from dotenv import load_dotenv
load_dotenv()
sys.path.insert(0, ".")

logging.basicConfig(level=logging.INFO, format="%(name)s - %(message)s")

from src.graph.state import AgentState
from src.agents.document_agent import document_agent_node

SOURCE = r"C:\Users\Batman\Downloads\22 CGI_Vendor_Final-2 (002).pptx"

state = AgentState(
    messages=[],
    intent="doc_task_pptx",
    source_doc_path=SOURCE,
    template_doc_path="",
    requirements=[
        "Analyze this ppt, I want you to follow the design theme of the slide 4 and apply it everywhere, font style, colour scheme and also check for spellings, fill the missing index slide stuff and make a new ppt with all these visual changes. Be vary of the margins and spaces and font sizes. do not make any mistakes and be thorough with your work."
    ],
    session_id="e2e_test",
    retry_count=0,
    critic_feedback="",
)

print("Running DocuForge document_agent_node...")
result = asyncio.run(document_agent_node(state))

print("\n=== Result ===")
print(f"output_doc_path: {result.get('output_doc_path')}")
print(f"requirements_met: {result.get('requirements_met')}")
print(f"draft_response: {result.get('draft_response', '')[:200]}")
print(f"retry_count: {result.get('retry_count')}")

if result.get('output_doc_path'):
    from pptx import Presentation
    p = Presentation(result['output_doc_path'])
    print(f"\nOutput slides: {len(p.slides)}")
    for i, slide in enumerate(p.slides, 1):
        shapes = [(s.name, s.shape_type) for s in slide.shapes]
        titles = [s.name for s in slide.shapes if s.shape_type == 14]  # PLACEHOLDERs
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        has_table = any(s.shape_type == MSO_SHAPE_TYPE.TABLE for s in slide.shapes)
        has_group = any(s.shape_type == MSO_SHAPE_TYPE.GROUP for s in slide.shapes)
        has_pic   = any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)
        print(f"  Slide {i}: {len(shapes)} shapes table={has_table} group={has_group} pic={has_pic}")
